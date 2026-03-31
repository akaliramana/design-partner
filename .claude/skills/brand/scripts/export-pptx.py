#!/usr/bin/env python3
"""
Export HTML deck to PPTX format using python-pptx.

Reads design tokens from MASTER.md, creates a branded PPTX with
matching colors and typography. HTML is the design source of truth;
PPTX is the practical export for sharing.

Usage:
    python3 export-pptx.py --master design-system/MASTER.md --title "Deck Title" --output output.pptx

Dependencies:
    pip install python-pptx Pillow
"""

import argparse
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def hex_to_rgb(hex_color: str) -> tuple:
    """Convert #RRGGBB to (R, G, B) tuple."""
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def parse_master_tokens(master_path: str) -> dict:
    """Extract brand tokens from MASTER.md."""
    content = Path(master_path).read_text(encoding="utf-8")
    tokens = {}

    # Extract colors
    for pattern, key in [
        (r"Primary[^\n]*`(#[0-9A-Fa-f]{6})`", "primary"),
        (r"Background[^\n]*`(#[0-9A-Fa-f]{6})`", "background"),
        (r"Text[^\n]*`(#[0-9A-Fa-f]{6})`", "text"),
        (r"CTA|Accent[^\n]*`(#[0-9A-Fa-f]{6})`", "accent"),
        (r"Secondary[^\n]*`(#[0-9A-Fa-f]{6})`", "secondary"),
    ]:
        match = re.search(pattern, content, re.IGNORECASE)
        if match:
            tokens[key] = match.group(1)

    # Extract fonts
    heading = re.search(r"Heading\s*(?:Font)?[:\s]*\**\s*([A-Za-z\s]+)", content)
    body = re.search(r"Body\s*(?:Font)?[:\s]*\**\s*([A-Za-z\s]+)", content)
    tokens["font_heading"] = heading.group(1).strip().rstrip("*") if heading else "Inter"
    tokens["font_body"] = body.group(1).strip().rstrip("*") if body else "Inter"

    return tokens


def create_branded_pptx(tokens: dict, title: str, output_path: str):
    """Create a branded PPTX template with custom colors and fonts."""
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 16:9
    prs.slide_height = Emu(6858000)

    bg_color = hex_to_rgb(tokens.get("background", "#0A0A0B"))
    primary_color = hex_to_rgb(tokens.get("primary", "#7C5CFC"))
    text_color = hex_to_rgb(tokens.get("text", "#E8E8ED"))

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_color)

    # Accent bar
    left = Inches(1.0)
    top = Inches(3.0)
    width = Inches(0.6)
    height = Inches(0.04)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*primary_color)
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(3.2), Inches(10.0), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*text_color)
    p.font.name = tokens.get("font_heading", "Inter")

    # Blank content slides (user fills in)
    for i in range(4):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg_color)

        # Slide number
        txBox = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.0), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 2)
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(100, 100, 110)
        p.alignment = PP_ALIGN.RIGHT

    prs.save(output_path)
    print(f"Created {output_path} ({len(prs.slides)} slides)")
    print(f"  Background: {tokens.get('background', '#0A0A0B')}")
    print(f"  Primary: {tokens.get('primary', '#7C5CFC')}")
    print(f"  Font: {tokens.get('font_heading', 'Inter')}")
    print(f"\nOpen in PowerPoint and add slide content. Brand colors and fonts are pre-applied.")


def main():
    if not HAS_PPTX:
        print("Error: python-pptx not installed.")
        print("Install with: pip install python-pptx Pillow")
        return

    parser = argparse.ArgumentParser(description="Export branded PPTX from MASTER.md tokens")
    parser.add_argument("--master", required=True, help="Path to design-system/MASTER.md")
    parser.add_argument("--title", default="Untitled Presentation", help="Deck title")
    parser.add_argument("--output", default="templates/deck-master.pptx", help="Output PPTX path")
    args = parser.parse_args()

    if not Path(args.master).exists():
        print(f"Error: {args.master} not found.")
        return

    tokens = parse_master_tokens(args.master)
    Path(args.output).parent.mkdir(parents=True, exist_ok=True)
    create_branded_pptx(tokens, args.title, args.output)


if __name__ == "__main__":
    main()
